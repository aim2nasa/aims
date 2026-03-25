"""
Document Metadata Extraction Service
"""
import os
import hashlib
import mimetypes
import logging
import io
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Optional, Tuple
import aiofiles

logger = logging.getLogger(__name__)

# Optional imports for enhanced functionality
try:
    import pypdfium2 as pdfium
    HAS_PYPDFIUM2 = True
except ImportError:
    HAS_PYPDFIUM2 = False
    logger.warning("pypdfium2 not available. PDF text extraction will be limited.")

try:
    import exifread
    HAS_EXIFREAD = True
except ImportError:
    HAS_EXIFREAD = False
    logger.warning("exifread not available. EXIF extraction will be limited.")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    logger.warning("Pillow not available. Image dimension extraction will be limited.")

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    logger.warning("openpyxl not available. Excel text extraction will be limited.")

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx not available. DOCX text extraction will be limited.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx not available. PPTX text extraction will be limited.")

try:
    import xlrd
    HAS_XLRD = True
except ImportError:
    HAS_XLRD = False
    logger.warning("xlrd not available. XLS text extraction will be limited.")


class MetaService:
    """Extract metadata from documents using class methods"""

    @classmethod
    async def extract_metadata(
        cls,
        file_path: Optional[str] = None,
        file_content: Optional[bytes] = None,
        filename: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Extract metadata from a file.

        Args:
            file_path: Path to the file on disk
            file_content: Binary content of the file
            filename: Original filename (used when file_content is provided)

        Returns:
            Metadata dictionary
        """
        try:
            if file_path:
                path = Path(file_path)
                if not path.exists():
                    return cls._error_response("FILE_NOT_FOUND", f"File not found: {file_path}")

                content = await cls._read_file(file_path)
                name = path.name
                stat = path.stat()
                size = stat.st_size
                created_at = datetime.fromtimestamp(stat.st_ctime).isoformat()
            elif file_content and filename:
                content = file_content
                name = filename
                size = len(file_content)
                created_at = datetime.utcnow().isoformat()
            else:
                return cls._error_response("NO_INPUT", "No file path or content provided")

            # Basic metadata
            extension = Path(name).suffix.lower()
            # HWP/HWPX는 OS별로 비표준 MIME을 반환하므로 확장자 기반으로 우선 보정
            _hwp_mime = {".hwp": "application/x-hwp", ".hwpx": "application/vnd.hancom.hwpx"}
            if extension in _hwp_mime:
                mime_type = _hwp_mime[extension]
            else:
                mime_type, _ = mimetypes.guess_type(name)
                if not mime_type:
                    mime_type = "application/octet-stream"

            # Compute file hash
            file_hash = hashlib.sha256(content).hexdigest()

            result = {
                "status": "OK",
                "filename": name,
                "extension": extension,
                "mime_type": mime_type,
                "file_size": size,
                "created_at": created_at,
                "file_hash": file_hash,
                "extracted_text": None,
                "num_pages": None,
                "pdf_text_ratio": None,
                "exif": None,
                "width": None,
                "height": None,
                "error": False
            }

            # Extract info based on file type
            if mime_type == "application/pdf":
                pdf_info = await cls._extract_pdf_info(content)
                result.update(pdf_info)
            elif mime_type and mime_type.startswith("image/"):
                image_info = cls._extract_image_info(content, mime_type)
                result.update(image_info)
            elif mime_type and mime_type.startswith("text/"):
                result["extracted_text"] = content.decode("utf-8", errors="ignore")
            # Office documents - XLSX
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel"
            ):
                xlsx_info = cls._extract_xlsx_info(content, mime_type)
                result.update(xlsx_info)
            # Office documents - DOCX
            elif mime_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword"
            ):
                docx_info = cls._extract_docx_info(content, mime_type)
                result.update(docx_info)
            # Office documents - PPTX
            elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                pptx_info = cls._extract_pptx_info(content)
                result.update(pptx_info)

            return result

        except Exception as e:
            logger.error(f"Metadata extraction error: {e}")
            return cls._error_response("EXTRACTION_ERROR", str(e))

    @classmethod
    async def _read_file(cls, file_path: str) -> bytes:
        """Read file content"""
        async with aiofiles.open(file_path, 'rb') as f:
            return await f.read()

    @classmethod
    async def _extract_pdf_info(cls, content: bytes) -> Dict[str, Any]:
        """Extract PDF-specific information"""
        result = {
            "num_pages": None,
            "extracted_text": None,
            "pdf_text_ratio": None
        }

        if not HAS_PYPDFIUM2:
            return result

        try:
            pdf = pdfium.PdfDocument(content)
            result["num_pages"] = len(pdf)

            # Extract text from all pages
            text_parts = []
            try:
                for page in pdf:
                    try:
                        textpage = page.get_textpage()
                        try:
                            text_parts.append(textpage.get_text_bounded())
                        finally:
                            textpage.close()
                    finally:
                        page.close()
            finally:
                pdf.close()

            full_text = "\n".join(text_parts)
            result["extracted_text"] = full_text

            # Calculate text ratio (characters per page)
            if result["num_pages"] > 0:
                result["pdf_text_ratio"] = len(full_text) / result["num_pages"]

        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            result["error"] = str(e)

        return result

    @classmethod
    def _extract_image_info(cls, content: bytes, mime_type: str) -> Dict[str, Any]:
        """Extract image-specific information including EXIF and dimensions"""
        result = {
            "exif": None,
            "width": None,
            "height": None
        }

        # Extract image dimensions using PIL
        if HAS_PIL:
            try:
                img = Image.open(io.BytesIO(content))
                result["width"] = img.width
                result["height"] = img.height
                img.close()
            except Exception as e:
                logger.warning(f"PIL image dimension extraction failed: {e}")

        # Extract EXIF data (primarily for JPEG)
        if HAS_EXIFREAD and mime_type in ("image/jpeg", "image/tiff"):
            try:
                tags = exifread.process_file(io.BytesIO(content), details=False)
                if tags:
                    exif_data = {}
                    for tag, value in tags.items():
                        # Skip thumbnail data and internal tags
                        if tag.startswith("Thumbnail") or tag.startswith("EXIF MakerNote"):
                            continue
                        # Convert IfdTag to string
                        str_value = str(value)
                        # Skip very long values (binary data)
                        if len(str_value) <= 500:
                            exif_data[tag] = str_value

                    if exif_data:
                        result["exif"] = exif_data

                        # Extract commonly used EXIF fields as top-level properties
                        if "EXIF DateTimeOriginal" in exif_data:
                            result["date_taken"] = exif_data["EXIF DateTimeOriginal"]
                        elif "EXIF DateTimeDigitized" in exif_data:
                            result["date_taken"] = exif_data["EXIF DateTimeDigitized"]
                        elif "Image DateTime" in exif_data:
                            result["date_taken"] = exif_data["Image DateTime"]

                        if "Image Make" in exif_data:
                            result["camera_make"] = exif_data["Image Make"]
                        if "Image Model" in exif_data:
                            result["camera_model"] = exif_data["Image Model"]

                        # GPS coordinates
                        if "GPS GPSLatitude" in exif_data and "GPS GPSLongitude" in exif_data:
                            result["gps_latitude"] = exif_data["GPS GPSLatitude"]
                            result["gps_longitude"] = exif_data["GPS GPSLongitude"]
                            if "GPS GPSLatitudeRef" in exif_data:
                                result["gps_latitude_ref"] = exif_data["GPS GPSLatitudeRef"]
                            if "GPS GPSLongitudeRef" in exif_data:
                                result["gps_longitude_ref"] = exif_data["GPS GPSLongitudeRef"]

                        # Image orientation
                        if "Image Orientation" in exif_data:
                            result["orientation"] = exif_data["Image Orientation"]

            except Exception as e:
                logger.warning(f"EXIF extraction failed: {e}")

        return result

    @classmethod
    def _extract_xlsx_info(cls, content: bytes, mime_type: str) -> Dict[str, Any]:
        """Extract text from Excel files (XLSX/XLS)"""
        result = {"extracted_text": None}

        # XLSX format (OpenXML)
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            if not HAS_OPENPYXL:
                return result
            try:
                wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                all_text = ""
                for sheet_name in wb.sheetnames:
                    all_text += f"\n--- 시트: {sheet_name} ---\n"
                    ws = wb[sheet_name]
                    for row in ws.iter_rows(values_only=True):
                        row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                        if row_text.strip():
                            all_text += row_text + "\n"
                wb.close()

                # Check if meaningful text exists
                cleaned = all_text.replace("\n", "").replace("\t", "").replace("--- 시트:", "").strip()
                if cleaned:
                    result["extracted_text"] = all_text
            except Exception as e:
                logger.warning(f"XLSX text extraction failed: {e}")

        # XLS format (Legacy)
        elif mime_type == "application/vnd.ms-excel":
            if not HAS_XLRD:
                return result
            try:
                wb = xlrd.open_workbook(file_contents=content)
                all_text = ""
                for sheet_name in wb.sheet_names():
                    all_text += f"\n--- 시트: {sheet_name} ---\n"
                    ws = wb.sheet_by_name(sheet_name)
                    for row_idx in range(ws.nrows):
                        row_text = "\t".join(str(ws.cell_value(row_idx, col_idx)) for col_idx in range(ws.ncols))
                        if row_text.strip():
                            all_text += row_text + "\n"

                cleaned = all_text.replace("\n", "").replace("\t", "").replace("--- 시트:", "").strip()
                if cleaned:
                    result["extracted_text"] = all_text
            except Exception as e:
                logger.warning(f"XLS text extraction failed: {e}")

        return result

    @classmethod
    def _extract_docx_info(cls, content: bytes, mime_type: str) -> Dict[str, Any]:
        """Extract text from Word documents (DOCX)"""
        result = {"extracted_text": None}

        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            if not HAS_DOCX:
                return result
            try:
                doc = DocxDocument(io.BytesIO(content))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                if paragraphs:
                    result["extracted_text"] = "\n".join(paragraphs)
            except Exception as e:
                logger.warning(f"DOCX text extraction failed: {e}")

        # DOC format not supported by python-docx
        elif mime_type == "application/msword":
            logger.info("DOC format not supported for text extraction (use DOCX)")

        return result

    @classmethod
    def _extract_pptx_info(cls, content: bytes) -> Dict[str, Any]:
        """Extract text from PowerPoint presentations (PPTX)"""
        result = {"extracted_text": None, "num_pages": None}

        if not HAS_PPTX:
            return result

        try:
            prs = Presentation(io.BytesIO(content))
            result["num_pages"] = len(prs.slides)

            all_text = ""
            for slide_idx, slide in enumerate(prs.slides, 1):
                all_text += f"\n--- 슬라이드 {slide_idx} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text += shape.text + "\n"

            cleaned = all_text.replace("\n", "").replace("--- 슬라이드", "").strip()
            if cleaned:
                result["extracted_text"] = all_text
        except Exception as e:
            logger.warning(f"PPTX text extraction failed: {e}")

        return result

    @classmethod
    def _error_response(cls, error_code: str, message: str) -> Dict[str, Any]:
        """Generate error response"""
        return {
            "status": 500,
            "error": True,
            "code": error_code,
            "message": message,
            "filename": None,
            "extension": None,
            "mime_type": None,
            "file_size": None,
            "created_at": None,
            "file_hash": None,
            "extracted_text": None,
            "num_pages": None,
            "pdf_text_ratio": None,
            "exif": None,
            "width": None,
            "height": None
        }
