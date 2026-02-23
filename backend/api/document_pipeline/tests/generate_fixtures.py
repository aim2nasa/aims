"""
Smoke Test Fixture Generator
=============================
tests/fixtures/ 디렉토리에 스모크 테스트용 샘플 파일 10종을 생성한다.

실행: python tests/generate_fixtures.py
결과: tests/fixtures/ 에 sample.pdf, sample.docx, ... 생성

경로2 파일(HWP, DOC, PPT)은 서버에서 LibreOffice 변환이 필요하므로,
이 스크립트에서는 DOCX/PPTX 소스만 생성하고 변환은 별도 수행한다.
"""
import os
import sys
from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"
KEYWORD = "AIMS_SMOKE_TEST"
DUMMY_TEXT = f"AIMS 스모크 테스트 문서\n{KEYWORD}\n이 문서는 배포 후 텍스트 추출 경로를 검증하기 위한 샘플입니다."
DUMMY_TEXT_EN = f"AIMS Smoke Test Document\n{KEYWORD}\nThis document verifies text extraction paths after deployment."


def ensure_dir():
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    print(f"[OK] fixtures dir: {FIXTURES_DIR}")


def generate_pdf():
    """경로1: 텍스트 PDF (PyMuPDF로 추출)"""
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)  # A4
    # PyMuPDF는 CJK 내장 지원 — 폰트 파일 임베딩 불필요
    text = f"{DUMMY_TEXT}\n\n{DUMMY_TEXT_EN}"
    page.insert_text((50, 50), text, fontsize=11, fontname="helv")

    out = FIXTURES_DIR / "sample.pdf"
    doc.save(str(out), garbage=4, deflate=True)
    doc.close()
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_docx():
    """경로1: DOCX (python-docx로 추출)"""
    from docx import Document

    doc = Document()
    doc.add_paragraph(DUMMY_TEXT)
    doc.add_paragraph(DUMMY_TEXT_EN)

    out = FIXTURES_DIR / "sample.docx"
    doc.save(str(out))
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_xlsx():
    """경로1: XLSX (openpyxl로 추출)"""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Smoke Test"
    ws["A1"] = "AIMS 스모크 테스트 문서"
    ws["A2"] = KEYWORD
    ws["A3"] = "이 문서는 배포 후 텍스트 추출 경로를 검증하기 위한 샘플입니다."
    ws["B1"] = "AIMS Smoke Test Document"
    ws["B2"] = KEYWORD

    out = FIXTURES_DIR / "sample.xlsx"
    wb.save(str(out))
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_pptx():
    """경로1: PPTX (python-pptx로 추출)"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout (최소 크기)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
    tf = txBox.text_frame
    tf.text = f"AIMS Smoke Test\n{KEYWORD}"

    out = FIXTURES_DIR / "sample.pptx"
    prs.save(str(out))
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_rtf():
    """경로2: RTF (LibreOffice → PDF 변환 후 추출)
    사용자 제공 RTF 파일의 구조를 유지하되 내용을 더미로 교체.
    """
    rtf_content = r"""{\rtf1\ansi\deff0
{\fonttbl{\f0\froman Times New Roman;}}
{\colortbl;\red0\green0\blue0;}
\pard\plain\f0\fs24
AIMS \uc0\u49828 \u47784 \u53356  \u53580 \u49828 \u53944  \u47928 \u49436 \par
AIMS_SMOKE_TEST\par
\u51060  \u47928 \u49436 \u45716  \u48176 \u54252  \u54980  \u53581 \u49828 \u53944  \u52628 \u52636  \u44221 \u47196 \u47484  \u44160 \u51613 \u54616 \u44592  \u50948 \u54620  \u49368 \u54540 \u51077 \u45768 \u45796 .\par
AIMS Smoke Test Document\par
AIMS_SMOKE_TEST\par
This document verifies text extraction paths after deployment.\par
}"""
    out = FIXTURES_DIR / "sample.rtf"
    out.write_text(rtf_content, encoding="ascii")
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_scan_pdf():
    """경로3: 스캔 PDF (텍스트 레이어 없음, OCR 필요)"""
    from PIL import Image, ImageDraw
    import fitz  # PyMuPDF
    import io

    # 작은 이미지 생성 (400x200)
    img = Image.new("RGB", (400, 200), "white")
    draw = ImageDraw.Draw(img)
    font = _get_pil_font(20)
    draw.text((20, 30), KEYWORD, fill="black", font=font)
    draw.text((20, 80), "Scanned document for OCR", fill="black", font=_get_pil_font(14))

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=70)

    # 이미지를 PDF로 (텍스트 레이어 없음)
    doc = fitz.open()
    page = doc.new_page(width=400, height=200)
    page.insert_image(page.rect, stream=buf.getvalue())

    out = FIXTURES_DIR / "sample_scan.pdf"
    doc.save(str(out), garbage=4, deflate=True)
    doc.close()
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


def generate_jpg():
    """경로3: JPG 이미지 (OCR 필요)"""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (400, 200), "white")
    draw = ImageDraw.Draw(img)
    font = _get_pil_font(20)
    draw.text((20, 30), KEYWORD, fill="black", font=font)
    draw.text((20, 80), "Image for OCR testing", fill="black", font=_get_pil_font(14))

    out = FIXTURES_DIR / "sample.jpg"
    img.save(str(out), "JPEG", quality=70)
    print(f"[OK] {out.name} ({out.stat().st_size:,} bytes)")


# --- Helpers ---

def _find_korean_font():
    """시스템에서 한글 폰트 경로를 찾는다."""
    candidates = [
        "C:/Windows/Fonts/malgun.ttf",       # 맑은 고딕
        "C:/Windows/Fonts/gulim.ttc",         # 굴림
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",  # Linux
        "/usr/share/fonts/nanum/NanumGothic.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


def _get_pil_font(size):
    """PIL용 폰트를 반환한다."""
    from PIL import ImageFont

    candidates = [
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/malgun.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def main():
    print("=== Smoke Test Fixture Generator ===\n")
    ensure_dir()

    generators = [
        ("PDF (Path 1)", generate_pdf),
        ("DOCX (Path 1)", generate_docx),
        ("XLSX (Path 1)", generate_xlsx),
        ("PPTX (Path 1)", generate_pptx),
        ("RTF (Path 2)", generate_rtf),
        ("Scan PDF (Path 3)", generate_scan_pdf),
        ("JPG (Path 3)", generate_jpg),
    ]

    success = 0
    for name, gen_func in generators:
        try:
            gen_func()
            success += 1
        except Exception as e:
            print(f"[FAIL] {name}: {e}")

    print(f"\n=== {success}/{len(generators)} files generated ===")
    print(f"\nNote: HWP, DOC, PPT는 서버에서 LibreOffice 변환 필요:")
    print(f"  soffice --headless --convert-to doc fixtures/sample.docx --outdir fixtures/")
    print(f"  soffice --headless --convert-to ppt fixtures/sample.pptx --outdir fixtures/")
    print(f"  soffice --headless --convert-to hwp fixtures/sample.docx --outdir fixtures/")

    return 0 if success == len(generators) else 1


if __name__ == "__main__":
    sys.exit(main())
